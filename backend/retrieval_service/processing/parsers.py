"""
Document parsing utilities for various file formats.

This module provides text extraction from multiple document formats including
PDF, DOCX, DOC, TXT, MD, PPTX, PPT, XLSX, and XLS.
"""

import io
import mimetypes
from typing import Optional

from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF for PDF
import openpyxl
from retrieval_service.infrastructure.logging import log_debug, log_info, log_warning, log_error

# Optional imports for legacy formats
try:
    import mammoth  # For .doc
except ImportError:
    mammoth = None


SUPPORTED_EXTS = {
    ".pdf", ".docx", ".doc", ".txt",
    ".pptx", ".ppt",
    ".xlsx", ".xls",
    ".md"
}


def isDOC(filename: str) -> bool:
    """
    Check whether a filename belongs to supported document types.
    
    Args:
        filename: Name of the file to check
        
    Returns:
        bool: True if filename has a supported document extension, False otherwise
    """
    if not isinstance(filename, str):
        return False

    lower = filename.lower()
    return any(lower.endswith(ext) for ext in SUPPORTED_EXTS)


# ========== Extractors ==========

def extract_text_from_pdf(data: bytes) -> str:
    """
    Extract text from PDF bytes using PyMuPDF.
    
    Args:
        data: Raw bytes of the PDF file
        
    Returns:
        str: Extracted text from all pages
    """
    text_chunks = []

    with fitz.open(stream=data, filetype="pdf") as doc:
        for page in doc:
            t = page.get_text()
            if t:
                text_chunks.append(t)

    return "\n".join(text_chunks).strip()


def extract_text_from_docx(data: bytes) -> str:
    """
    Extract text from DOCX bytes.
    
    Args:
        data: Raw bytes of the DOCX file
        
    Returns:
        str: Extracted text from all paragraphs
    """
    file_like = io.BytesIO(data)
    doc = Document(file_like)
    return "\n".join(
        para.text for para in doc.paragraphs if para.text.strip()
    )


def extract_text_from_doc(data: bytes) -> str:
    """
    Extract text from legacy .doc files.
    Uses mammoth if available, falls back to empty string otherwise.
    
    Args:
        data: Raw bytes of the DOC file
        
    Returns:
        str: Extracted text in markdown format, or error message if mammoth not available
    """
    if not mammoth:
        log_warning("[Parser] Mammoth library not available for .doc parsing")
        return "[.doc parsing not available — install `mammoth` to enable]"

    file_like = io.BytesIO(data)
    result = mammoth.convert_to_markdown(file_like)
    md_text = result.value  # markdown
    return md_text.strip()


def extract_text_from_txt(data: bytes) -> str:
    """
    Extract text from TXT bytes with UTF-8 fallback.
    
    Args:
        data: Raw bytes of the TXT file
        
    Returns:
        str: Decoded text content
    """
    try:
        return data.decode("utf-8").strip()
    except UnicodeDecodeError:
        log_debug("[Parser] UTF-8 decode failed, falling back to latin1")
        return data.decode("latin1").strip()


def extract_text_from_md(data: bytes) -> str:
    """
    Markdown treated as plain text.
    
    Args:
        data: Raw bytes of the MD file
        
    Returns:
        str: Decoded text content
    """
    return extract_text_from_txt(data)


def extract_text_from_pptx(data: bytes) -> str:
    """
    Extract text from PPTX bytes.
    
    Args:
        data: Raw bytes of the PPTX file
        
    Returns:
        str: Extracted text from all slides and shapes
    """
    file_like = io.BytesIO(data)
    prs = Presentation(file_like)

    text_chunks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_chunks.append(shape.text)

    return "\n".join(text_chunks).strip()


def extract_text_from_ppt(data: bytes) -> str:
    """
    Legacy .ppt — very limited text extraction.
    
    Args:
        data: Raw bytes of the PPT file
        
    Returns:
        str: Error message indicating limited support
    """
    log_warning("[Parser] Legacy .ppt format not fully supported")
    return "[Legacy .ppt parsing not implemented — convert to .pptx]"


def extract_text_from_xlsx(data: bytes) -> str:
    """
    Extract text from XLSX bytes.
    
    Args:
        data: Raw bytes of the XLSX file
        
    Returns:
        str: Extracted text from all cells
    """
    file_like = io.BytesIO(data)
    wb = openpyxl.load_workbook(file_like, data_only=True)
    text_chunks = []

    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            for cell in row:
                if cell is not None:
                    text_chunks.append(str(cell))

    return "\n".join(text_chunks).strip()


def extract_text_from_xls(data: bytes) -> str:
    """
    Legacy XLS — limited support unless xlrd installed.
    
    Args:
        data: Raw bytes of the XLS file
        
    Returns:
        str: Extracted text from all cells, or error message if xlrd not available
    """
    try:
        import xlrd
    except ImportError:
        log_warning("[Parser] xlrd library not available for .xls parsing")
        return "[.xls parsing not available — install xlrd to enable]"

    file_like = io.BytesIO(data)
    book = xlrd.open_workbook(file_contents=data)
    text_chunks = []

    for sheet in book.sheets():
        for row_idx in range(sheet.nrows):
            cells = sheet.row_values(row_idx)
            text_chunks.extend(str(c) for c in cells if c)

    return "\n".join(text_chunks).strip()


# ========== Dispatcher ==========

def extractDOC(data: bytes, filename: Optional[str] = None, max_chars: int = 15000) -> str:
    """
    Extract text from many document formats.
    Automatically detects document type using filename extension.
    Applies length restriction to control LLM cost.
    
    Args:
        data: Raw bytes of the document file
        filename: Name of the file (required for type detection)
        max_chars: Maximum characters to return (default: 15000)
        
    Returns:
        str: Extracted text, truncated if necessary
        
    Raises:
        ValueError: If filename is not provided or format is unsupported
    """
    if not filename:
        raise ValueError("filename is required to detect document type")

    name = filename.lower()

    try:
        # Determine extractor
        if name.endswith(".pdf"):
            log_debug(f"[Parser] Extracting PDF: {filename}")
            text = extract_text_from_pdf(data)

        elif name.endswith(".docx"):
            log_debug(f"[Parser] Extracting DOCX: {filename}")
            text = extract_text_from_docx(data)

        elif name.endswith(".doc"):
            log_debug(f"[Parser] Extracting DOC: {filename}")
            text = extract_text_from_doc(data)

        elif name.endswith(".txt"):
            log_debug(f"[Parser] Extracting TXT: {filename}")
            text = extract_text_from_txt(data)

        elif name.endswith(".md"):
            log_debug(f"[Parser] Extracting MD: {filename}")
            text = extract_text_from_md(data)

        elif name.endswith(".pptx"):
            log_debug(f"[Parser] Extracting PPTX: {filename}")
            text = extract_text_from_pptx(data)

        elif name.endswith(".ppt"):
            log_debug(f"[Parser] Extracting PPT: {filename}")
            text = extract_text_from_ppt(data)

        elif name.endswith(".xlsx"):
            log_debug(f"[Parser] Extracting XLSX: {filename}")
            text = extract_text_from_xlsx(data)

        elif name.endswith(".xls"):
            log_debug(f"[Parser] Extracting XLS: {filename}")
            text = extract_text_from_xls(data)

        else:
            log_error(f"[Parser] Unsupported document format: {filename}")
            raise ValueError(f"Unsupported document format: {filename}")
    except Exception as e:
        log_error(f"[Parser] Error extracting text from {filename}: {e}")
        text = f"[Error extracting text from {filename}: {e}]"

    # Clean whitespace
    if text:
        text = text.replace("\r", "").strip()

    # Enforce length limit
    if len(text) > max_chars:
        log_debug(f"[Parser] Truncating text from {len(text)} to {max_chars} chars")
        text = text[:max_chars] + "\n...[text truncated]"

    return text
