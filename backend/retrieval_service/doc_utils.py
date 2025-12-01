# doc_utils.py
import io
import mimetypes
from typing import Optional

from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF for PDF
import openpyxl

# Optional imports for legacy formats
try:
    import mammoth  # For .doc
except ImportError:
    mammoth = None


SUPPORTED_EXTS = {
    ".pdf", ".docx", ".doc", ".txt",
    ".pptx", ".ppt",
    ".xlsx", ".xls",
    ".md"
}


def isDOC(filename: str) -> bool:
    """Check whether a filename belongs to supported document types."""
    if not isinstance(filename, str):
        return False

    lower = filename.lower()
    return any(lower.endswith(ext) for ext in SUPPORTED_EXTS)


# ========== Extractors ==========

def extract_text_from_pdf(data: bytes) -> str:
    """Extract text from PDF bytes using PyMuPDF."""
    text_chunks = []

    with fitz.open(stream=data, filetype="pdf") as doc:
        for page in doc:
            t = page.get_text()
            if t:
                text_chunks.append(t)

    return "\n".join(text_chunks).strip()


def extract_text_from_docx(data: bytes) -> str:
    """Extract text from DOCX bytes."""
    file_like = io.BytesIO(data)
    doc = Document(file_like)
    return "\n".join(
        para.text for para in doc.paragraphs if para.text.strip()
    )


def extract_text_from_doc(data: bytes) -> str:
    """
    Extract text from legacy .doc files.
    Uses mammoth if available, falls back to empty string otherwise.
    """
    if not mammoth:
        return "[.doc parsing not available — install `mammoth` to enable]"

    file_like = io.BytesIO(data)
    result = mammoth.convert_to_markdown(file_like)
    md_text = result.value  # markdown
    return md_text.strip()


def extract_text_from_txt(data: bytes) -> str:
    """Extract text from TXT bytes with UTF-8 fallback."""
    try:
        return data.decode("utf-8").strip()
    except UnicodeDecodeError:
        return data.decode("latin1").strip()


def extract_text_from_md(data: bytes) -> str:
    """Markdown treated as plain text."""
    return extract_text_from_txt(data)


def extract_text_from_pptx(data: bytes) -> str:
    """Extract text from PPTX bytes."""
    file_like = io.BytesIO(data)
    prs = Presentation(file_like)

    text_chunks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_chunks.append(shape.text)

    return "\n".join(text_chunks).strip()


def extract_text_from_ppt(data: bytes) -> str:
    """Legacy .ppt — very limited text extraction."""
    return "[Legacy .ppt parsing not implemented — convert to .pptx]"


def extract_text_from_xlsx(data: bytes) -> str:
    """Extract text from XLSX bytes."""
    file_like = io.BytesIO(data)
    wb = openpyxl.load_workbook(file_like, data_only=True)
    text_chunks = []

    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            for cell in row:
                if cell is not None:
                    text_chunks.append(str(cell))

    return "\n".join(text_chunks).strip()


def extract_text_from_xls(data: bytes) -> str:
    """Legacy XLS — limited support unless xlrd installed."""
    try:
        import xlrd
    except ImportError:
        return "[.xls parsing not available — install xlrd to enable]"

    file_like = io.BytesIO(data)
    book = xlrd.open_workbook(file_contents=data)
    text_chunks = []

    for sheet in book.sheets():
        for row_idx in range(sheet.nrows):
            cells = sheet.row_values(row_idx)
            text_chunks.extend(str(c) for c in cells if c)

    return "\n".join(text_chunks).strip()


# ========== Dispatcher ==========

def extractDOC(data: bytes, filename: Optional[str] = None, max_chars: int = 15000) -> str:
    """
    Extract text from many document formats.
    Automatically detects document type using filename extension.
    Applies length restriction to control LLM cost.
    """
    if not filename:
        raise ValueError("filename is required to detect document type")

    name = filename.lower()

    try:
        # Determine extractor
        if name.endswith(".pdf"):
            text = extract_text_from_pdf(data)

        elif name.endswith(".docx"):
            text = extract_text_from_docx(data)

        elif name.endswith(".doc"):
            text = extract_text_from_doc(data)

        elif name.endswith(".txt"):
            text = extract_text_from_txt(data)

        elif name.endswith(".md"):
            text = extract_text_from_md(data)

        elif name.endswith(".pptx"):
            text = extract_text_from_pptx(data)

        elif name.endswith(".ppt"):
            text = extract_text_from_ppt(data)

        elif name.endswith(".xlsx"):
            text = extract_text_from_xlsx(data)

        elif name.endswith(".xls"):
            text = extract_text_from_xls(data)

        else:
            raise ValueError(f"Unsupported document format: {filename}")
    except Exception as e:
        text = f"[Error extracting text from {filename}: {e}]"

    # Clean whitespace
    if text:
        text = text.replace("\r", "").strip()

    # Enforce length limit
    if len(text) > max_chars:
        text = text[:max_chars] + "\n...[text truncated]"

    return text
